"""
NS Artifact Renderer
Turns stabilized outputs into artifacts.

Three renderers:
    render_infographic_svg(spec) → SVG string
    render_deck_pptx(spec, output_path) → .pptx file
    render_brief_pdf(spec, output_path) → .pdf file

Spec format (shared):
{
    "title": str,
    "subtitle": str,
    "theme": "dark" | "light" | "axiolev",
    "sections": [
        {
            "heading": str,
            "body": str,
            "stats": [{"label": str, "value": str}],
            "bullets": [str],
        }
    ],
    "source_domain": str,     # which NS domain this came from
    "commit_receipt_id": str, # receipt of the commit that produced this
    "generated_at": str,
}

All outputs logged to /runs/ with content hash.
"""

import hashlib
import json
import time
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Dict, List, Optional


RUNS_DIR = Path("/tmp/ns_workspace/runs")
RUNS_DIR.mkdir(parents=True, exist_ok=True)

THEMES = {
    "axiolev": {
        "bg": "#0a0a12",
        "bg2": "#12121e",
        "accent": "#4a9eff",
        "accent2": "#7c3aed",
        "text": "#e8e8f0",
        "muted": "#6b7280",
        "border": "#1e1e2e",
        "gold": "#f59e0b",
    },
    "dark": {
        "bg": "#111827",
        "bg2": "#1f2937",
        "accent": "#3b82f6",
        "accent2": "#8b5cf6",
        "text": "#f9fafb",
        "muted": "#6b7280",
        "border": "#374151",
        "gold": "#f59e0b",
    },
    "light": {
        "bg": "#ffffff",
        "bg2": "#f3f4f6",
        "accent": "#2563eb",
        "accent2": "#7c3aed",
        "text": "#111827",
        "muted": "#6b7280",
        "border": "#e5e7eb",
        "gold": "#d97706",
    },
}

def _ts() -> str:
    return datetime.now(timezone.utc).isoformat()

def _log_run(kind: str, spec: dict, output_path: str) -> str:
    """Log artifact creation to /runs/ with content hash."""
    run_id = hashlib.sha256(f"{kind}{time.time_ns()}".encode()).hexdigest()[:12]
    log = {
        "run_id": run_id,
        "kind": kind,
        "title": spec.get("title", ""),
        "output_path": output_path,
        "spec_hash": hashlib.sha256(json.dumps(spec, sort_keys=True).encode()).hexdigest()[:16],
        "source_domain": spec.get("source_domain", ""),
        "commit_receipt_id": spec.get("commit_receipt_id", ""),
        "generated_at": _ts(),
    }
    log_path = RUNS_DIR / f"{run_id}_{kind}.json"
    log_path.write_text(json.dumps(log, indent=2))
    return run_id


# ── SVG Infographic ─────────────────────────────────────────────────────────

def render_infographic_svg(spec: dict) -> str:
    """
    Generate a professional SVG infographic from spec.
    Returns SVG string.
    """
    theme_name = spec.get("theme", "axiolev")
    c = THEMES.get(theme_name, THEMES["axiolev"])
    title = spec.get("title", "NS∞ Report")
    subtitle = spec.get("subtitle", "")
    sections = spec.get("sections", [])
    domain = spec.get("source_domain", "")

    # Canvas
    W = 900
    SECTION_H = 200
    HEADER_H = 120
    FOOTER_H = 50
    H = HEADER_H + max(len(sections), 1) * SECTION_H + FOOTER_H + 40

    parts = [f"""<svg xmlns="http://www.w3.org/2000/svg" width="{W}" height="{H}"
  viewBox="0 0 {W} {H}" font-family="system-ui,-apple-system,sans-serif">
  <defs>
    <linearGradient id="headerGrad" x1="0%" y1="0%" x2="100%" y2="100%">
      <stop offset="0%" style="stop-color:{c['accent']};stop-opacity:0.15"/>
      <stop offset="100%" style="stop-color:{c['accent2']};stop-opacity:0.05"/>
    </linearGradient>
    <linearGradient id="accentLine" x1="0%" y1="0%" x2="100%" y2="0%">
      <stop offset="0%" style="stop-color:{c['accent']}"/>
      <stop offset="100%" style="stop-color:{c['accent2']}"/>
    </linearGradient>
  </defs>

  <!-- Background -->
  <rect width="{W}" height="{H}" fill="{c['bg']}"/>

  <!-- Header area -->
  <rect width="{W}" height="{HEADER_H}" fill="url(#headerGrad)"/>
  <rect x="0" y="{HEADER_H - 2}" width="{W}" height="2" fill="url(#accentLine)"/>

  <!-- Logo mark -->
  <text x="32" y="38" font-size="11" fill="{c['accent']}" letter-spacing="3"
        font-weight="600">NS∞ · AXIOLEV</text>

  <!-- Title -->
  <text x="32" y="72" font-size="26" fill="{c['text']}" font-weight="700">{_esc(title)}</text>
  <text x="32" y="98" font-size="13" fill="{c['muted']}">{_esc(subtitle)}</text>

  <!-- Domain badge -->"""]

    if domain:
        parts.append(f"""  <rect x="{W - 150}" y="20" width="120" height="26"
        rx="13" fill="{c['accent']}22" stroke="{c['accent']}55" stroke-width="1"/>
  <text x="{W - 90}" y="37" text-anchor="middle" font-size="10"
        fill="{c['accent']}" letter-spacing="1">{_esc(domain.upper())}</text>""")

    # Sections
    for i, section in enumerate(sections):
        y = HEADER_H + 20 + i * SECTION_H
        heading = section.get("heading", "")
        body = section.get("body", "")
        stats = section.get("stats", [])
        bullets = section.get("bullets", [])

        # Section card
        parts.append(f"""
  <!-- Section {i+1} -->
  <rect x="24" y="{y}" width="{W - 48}" height="{SECTION_H - 16}"
        rx="8" fill="{c['bg2']}" stroke="{c['border']}" stroke-width="1"/>
  <rect x="24" y="{y}" width="4" height="{SECTION_H - 16}"
        rx="2" fill="{c['accent'] if i % 2 == 0 else c['accent2']}"/>

  <!-- Heading -->
  <text x="44" y="{y + 28}" font-size="14" fill="{c['text']}"
        font-weight="600">{_esc(heading)}</text>""")

        # Stats row
        if stats:
            stat_w = min(150, (W - 80) // len(stats))
            for j, stat in enumerate(stats[:5]):
                sx = 44 + j * (stat_w + 12)
                parts.append(f"""
  <rect x="{sx}" y="{y + 44}" width="{stat_w}" height="60"
        rx="6" fill="{c['bg']}" stroke="{c['accent']}33" stroke-width="1"/>
  <text x="{sx + stat_w//2}" y="{y + 76}" text-anchor="middle"
        font-size="22" fill="{c['gold']}" font-weight="700">{_esc(str(stat.get('value', '')))}</text>
  <text x="{sx + stat_w//2}" y="{y + 94}" text-anchor="middle"
        font-size="9" fill="{c['muted']}" letter-spacing="1">{_esc(stat.get('label', '').upper())}</text>""")
            body_y = y + 118
        else:
            body_y = y + 48

        # Body text (wrapped at ~80 chars per line, max 3 lines)
        if body:
            lines = _wrap_text(body, 90)[:3]
            for li, line in enumerate(lines):
                parts.append(f"""
  <text x="44" y="{body_y + li * 18}" font-size="12" fill="{c['muted']}">{_esc(line)}</text>""")

        # Bullets
        if bullets and not body:
            for bi, bullet in enumerate(bullets[:4]):
                by = body_y + bi * 22
                parts.append(f"""
  <circle cx="52" cy="{by - 4}" r="3" fill="{c['accent']}"/>
  <text x="62" y="{by}" font-size="12" fill="{c['text']}">{_esc(bullet[:80])}</text>""")

    # Footer
    footer_y = HEADER_H + 20 + max(len(sections), 1) * SECTION_H + 16
    ts = datetime.now(timezone.utc).strftime("%Y-%m-%d %H:%M UTC")
    parts.append(f"""
  <!-- Footer -->
  <rect x="0" y="{footer_y}" width="{W}" height="1" fill="{c['border']}"/>
  <text x="32" y="{footer_y + 28}" font-size="9" fill="{c['muted']}"
        letter-spacing="1">GENERATED BY NS∞ · AXIOLEV HOLDINGS · {_esc(ts)}</text>
  <text x="{W - 32}" y="{footer_y + 28}" text-anchor="end" font-size="9"
        fill="{c['accent']}">northstar.axiolev.com</text>

</svg>""")

    svg = "\n".join(parts)
    _log_run("infographic_svg", spec, "in-memory")
    return svg


# ── PPTX Deck ───────────────────────────────────────────────────────────────

def render_deck_pptx(spec: dict, output_path: str) -> str:
    """
    Generate a PPTX presentation from spec.
    Returns output_path.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    def _hex(h: str) -> RGBColor:
        h = h.lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    theme_name = spec.get("theme", "axiolev")
    c = THEMES.get(theme_name, THEMES["axiolev"])
    title_text = spec.get("title", "NS∞ Report")
    subtitle_text = spec.get("subtitle", "")
    sections = spec.get("sections", [])
    domain = spec.get("source_domain", "")

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    SL = prs.slide_layouts[6]  # blank

    def _add_rect(slide, l, t, w, h, fill_hex, alpha=None):
        from pptx.util import Pt
        shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
        shape.line.fill.background()
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex(fill_hex)
        return shape

    def _add_text(slide, text, l, t, w, h, size, color_hex,
                  bold=False, align="left"):
        txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = txb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        if align == "center":
            p.alignment = PP_ALIGN.CENTER
        elif align == "right":
            p.alignment = PP_ALIGN.RIGHT
        run = p.runs[0]
        run.font.size = Pt(size)
        run.font.color.rgb = _hex(color_hex)
        run.font.bold = bold
        return txb

    # ── Slide 1: Title ──────────────────────────────────────────────────────
    slide = prs.slides.add_slide(SL)
    _add_rect(slide, 0, 0, 13.33, 7.5, c["bg"].lstrip("#"))
    _add_rect(slide, 0, 0, 13.33, 0.05, c["accent"].lstrip("#"))
    _add_rect(slide, 0, 7.45, 13.33, 0.05, c["accent2"].lstrip("#"))

    # NS∞ mark
    _add_text(slide, "NS∞ · AXIOLEV", 0.5, 0.3, 4, 0.4, 10,
              c["accent"].lstrip("#"))

    # Big title
    _add_text(slide, title_text, 0.5, 2.2, 10, 1.2, 44,
              c["text"].lstrip("#"), bold=True)

    # Subtitle
    _add_text(slide, subtitle_text, 0.5, 3.6, 10, 0.6, 18,
              c["muted"].lstrip("#"))

    # Domain badge
    if domain:
        _add_rect(slide, 9.5, 6.7, 3.3, 0.5, c["bg2"].lstrip("#"))
        _add_text(slide, domain.upper(), 9.5, 6.7, 3.3, 0.5, 11,
                  c["accent"].lstrip("#"), align="center")

    # Date
    ts = datetime.now(timezone.utc).strftime("%B %d, %Y")
    _add_text(slide, ts, 0.5, 6.9, 5, 0.4, 10, c["muted"].lstrip("#"))

    # ── Section Slides ───────────────────────────────────────────────────────
    for i, section in enumerate(sections):
        slide = prs.slides.add_slide(SL)
        _add_rect(slide, 0, 0, 13.33, 7.5, c["bg"].lstrip("#"))
        _add_rect(slide, 0, 0, 0.06, 7.5,
                  (c["accent"] if i % 2 == 0 else c["accent2"]).lstrip("#"))
        _add_rect(slide, 0, 0, 13.33, 1.2, c["bg2"].lstrip("#"))

        heading = section.get("heading", "")
        body = section.get("body", "")
        stats = section.get("stats", [])
        bullets = section.get("bullets", [])

        _add_text(slide, heading, 0.3, 0.25, 11, 0.7, 28,
                  c["text"].lstrip("#"), bold=True)

        # Section number
        _add_text(slide, f"{i+1:02d}", 11.8, 0.3, 1.3, 0.6, 22,
                  c["accent"].lstrip("#"), bold=True, align="right")

        if stats:
            stat_w = min(2.5, 12.0 / len(stats))
            for j, stat in enumerate(stats[:4]):
                sx = 0.3 + j * (stat_w + 0.1)
                _add_rect(slide, sx, 1.4, stat_w, 1.8, c["bg2"].lstrip("#"))
                _add_text(slide, str(stat.get("value", "")),
                          sx, 1.7, stat_w, 0.9, 32,
                          c["gold"].lstrip("#"), bold=True, align="center")
                _add_text(slide, stat.get("label", ""),
                          sx, 2.7, stat_w, 0.4, 10,
                          c["muted"].lstrip("#"), align="center")

            if body:
                _add_text(slide, body, 0.3, 3.5, 12.5, 3.5, 14,
                          c["muted"].lstrip("#"))
        else:
            content_y = 1.4
            if body:
                _add_text(slide, body, 0.3, content_y, 12.5, 2.0, 16,
                          c["text"].lstrip("#"))
                content_y += 2.2
            if bullets:
                for bi, bullet in enumerate(bullets[:6]):
                    _add_text(slide, f"◆  {bullet}", 0.5,
                              content_y + bi * 0.7, 12, 0.65, 14,
                              c["text"].lstrip("#"))

        # Footer
        _add_rect(slide, 0, 7.2, 13.33, 0.01, c["border"].lstrip("#"))
        _add_text(slide, "NS∞ · AXIOLEV HOLDINGS", 0.3, 7.2, 6, 0.3,
                  8, c["muted"].lstrip("#"))
        _add_text(slide, f"{i+2}", 12.8, 7.2, 0.5, 0.3, 8,
                  c["muted"].lstrip("#"), align="right")

    # ── Closing slide ───────────────────────────────────────────────────────
    slide = prs.slides.add_slide(SL)
    _add_rect(slide, 0, 0, 13.33, 7.5, c["bg"].lstrip("#"))
    _add_rect(slide, 0, 3.5, 13.33, 0.04,
              c["accent"].lstrip("#"))
    _add_text(slide, "NS∞", 0.5, 2.5, 12.33, 1.2, 72,
              c["accent"].lstrip("#"), bold=True, align="center")
    _add_text(slide, "AXIOLEV HOLDINGS", 0.5, 3.8, 12.33, 0.6, 14,
              c["muted"].lstrip("#"), align="center")

    prs.save(output_path)
    _log_run("deck_pptx", spec, output_path)
    return output_path


# ── PDF Brief ────────────────────────────────────────────────────────────────

def render_brief_pdf(spec: dict, output_path: str) -> str:
    """
    Generate a PDF brief from spec.
    Returns output_path.
    """
    from reportlab.lib.pagesizes import letter
    from reportlab.lib import colors
    from reportlab.lib.units import inch
    from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer,
                                    Table, TableStyle, HRFlowable)
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_RIGHT

    theme_name = spec.get("theme", "axiolev")
    c = THEMES.get(theme_name, THEMES["axiolev"])

    def _rgb(h: str):
        h = h.lstrip("#")
        return colors.HexColor(f"#{h}")

    doc = SimpleDocTemplate(
        output_path,
        pagesize=letter,
        leftMargin=0.75*inch,
        rightMargin=0.75*inch,
        topMargin=0.75*inch,
        bottomMargin=0.75*inch,
    )

    styles = getSampleStyleSheet()

    style_title = ParagraphStyle("NSTitle",
        fontSize=24, fontName="Helvetica-Bold",
        textColor=_rgb(c["text"]), spaceAfter=6, leading=28)

    style_subtitle = ParagraphStyle("NSSubtitle",
        fontSize=12, fontName="Helvetica",
        textColor=_rgb(c["muted"]), spaceAfter=18)

    style_heading = ParagraphStyle("NSHeading",
        fontSize=16, fontName="Helvetica-Bold",
        textColor=_rgb(c["accent"]), spaceBefore=18, spaceAfter=6)

    style_body = ParagraphStyle("NSBody",
        fontSize=11, fontName="Helvetica",
        textColor=_rgb(c["text"]), spaceAfter=8, leading=16)

    style_bullet = ParagraphStyle("NSBullet",
        fontSize=11, fontName="Helvetica",
        textColor=_rgb(c["text"]), leftIndent=16,
        spaceAfter=4, leading=15, bulletIndent=4)

    style_meta = ParagraphStyle("NSMeta",
        fontSize=8, fontName="Helvetica",
        textColor=_rgb(c["muted"]), spaceAfter=4)

    story = []

    # Header
    title = spec.get("title", "NS∞ Report")
    subtitle = spec.get("subtitle", "")
    domain = spec.get("source_domain", "")
    ts = datetime.now(timezone.utc).strftime("%B %d, %Y · %H:%M UTC")
    commit_id = spec.get("commit_receipt_id", "")

    story.append(Paragraph("NS∞ · AXIOLEV HOLDINGS", style_meta))
    story.append(HRFlowable(width="100%", thickness=2,
                             color=_rgb(c["accent"]), spaceAfter=12))
    story.append(Paragraph(title, style_title))
    if subtitle:
        story.append(Paragraph(subtitle, style_subtitle))

    meta_parts = [ts]
    if domain:
        meta_parts.append(f"Domain: {domain}")
    if commit_id:
        meta_parts.append(f"Commit: {commit_id[:12]}")
    story.append(Paragraph(" · ".join(meta_parts), style_meta))
    story.append(HRFlowable(width="100%", thickness=0.5,
                             color=_rgb(c["border"]), spaceAfter=20))

    # Sections
    for i, section in enumerate(spec.get("sections", [])):
        heading = section.get("heading", "")
        body = section.get("body", "")
        stats = section.get("stats", [])
        bullets = section.get("bullets", [])

        story.append(Paragraph(heading, style_heading))

        if stats:
            cols = len(stats)
            col_w = (6.5 * inch) / cols
            stat_data = [[s.get("value", "") for s in stats],
                         [s.get("label", "").upper() for s in stats]]
            tbl = Table(stat_data, colWidths=[col_w] * cols)
            tbl.setStyle(TableStyle([
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("FONTSIZE", (0, 0), (-1, 0), 20),
                ("FONTNAME", (0, 1), (-1, 1), "Helvetica"),
                ("FONTSIZE", (0, 1), (-1, 1), 8),
                ("TEXTCOLOR", (0, 0), (-1, 0), _rgb(c["gold"])),
                ("TEXTCOLOR", (0, 1), (-1, 1), _rgb(c["muted"])),
                ("ALIGN", (0, 0), (-1, -1), "CENTER"),
                ("BACKGROUND", (0, 0), (-1, -1), _rgb(c["bg2"])),
                ("ROWBACKGROUNDS", (0, 0), (-1, -1), [_rgb(c["bg2"])]),
                ("BOX", (0, 0), (-1, -1), 0.5, _rgb(c["border"])),
                ("INNERGRID", (0, 0), (-1, -1), 0.25, _rgb(c["border"])),
                ("TOPPADDING", (0, 0), (-1, -1), 10),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 10),
            ]))
            story.append(tbl)
            story.append(Spacer(1, 0.15*inch))

        if body:
            story.append(Paragraph(body, style_body))

        for bullet in bullets:
            story.append(Paragraph(f"• {bullet}", style_bullet))

        story.append(HRFlowable(width="100%", thickness=0.25,
                                 color=_rgb(c["border"]), spaceAfter=8))

    # Footer
    story.append(Spacer(1, 0.2*inch))
    story.append(Paragraph(
        f"Generated by NS∞ Constitutional AI · AXIOLEV Holdings LLC · {ts}",
        style_meta
    ))

    doc.build(story)
    _log_run("brief_pdf", spec, output_path)
    return output_path


# ── List recent runs ─────────────────────────────────────────────────────────

def list_runs(limit: int = 20) -> List[dict]:
    runs = []
    for f in sorted(RUNS_DIR.glob("*.json"), key=lambda x: x.stat().st_mtime, reverse=True)[:limit]:
        try:
            runs.append(json.loads(f.read_text()))
        except Exception:
            pass
    return runs


# ── Helpers ──────────────────────────────────────────────────────────────────

def _esc(s: str) -> str:
    """XML-escape for SVG text."""
    return (str(s)
            .replace("&", "&amp;")
            .replace("<", "&lt;")
            .replace(">", "&gt;")
            .replace('"', "&quot;"))

def _wrap_text(text: str, width: int) -> List[str]:
    """Simple word-wrap."""
    words = text.split()
    lines, current = [], ""
    for word in words:
        if len(current) + len(word) + 1 <= width:
            current = (current + " " + word).strip()
        else:
            if current:
                lines.append(current)
            current = word
    if current:
        lines.append(current)
    return lines
